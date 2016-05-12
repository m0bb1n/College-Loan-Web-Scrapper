### GET IMAGE SCRAPER TO WORK
### SET UP DATA TABLE FOR POWER POINT

from pptx import Presentation #lets us create presentation
from pptx.util import Inches, Pt #lets us get pptx fonts, measurements etc
import time #lets us pause for a certain time
import requests # lets us download image from url
import urllib
import os # lets us execute console commands
from datetime import datetime  
from selenium import webdriver #lets us search the web
from selenium.webdriver.common.by import By #lets us search for elements
from selenium.webdriver.common.keys import Keys #lets us enter keys into webdriver
from selenium.webdriver.remote.webelement import WebElement 

class loan (object):
        simple = { }
        compound = { }
        collegeTuition = start.collegeTuition
        rate = 0.04

        def __init__(self, rate):
                self.rate = rate

        def find_simple (self):
                self.simple['equation'] = 'Interest = ' + str(self.collegeTuition) + '(' + str(self.rate) + '(year)'
                
                for year in range(1,11):
                        if year == 1 or year == 5 or year == 10:
                                full = self.collegeTuition * self.rate * year
                                self.simple[year] = str(full)

        def find_compound (self):
                self.compound['equation'] = 'Interest = ' + str(self.collegeTuition) + '(' + str(self.rate) + ')^' + 'year)'
                
                for year in range(1, 11):
                        if year == 1 or year == 5 or year == 10:
                                full = self.collegeTuition * ( (self.rate+1) ** year)
                                self.compound[year] = str(full)

class powerpoint (object):
	
	def __init__(self):

		###########################################################################
		#	 FIRST TIME RUN TO CREATE PRESENTATION AND MAKE DIRECTORY	  # 
		###########################################################################

		os.system('mkdir -p ~/Desktop/College_Loan_Presentations/' + start.collegeDir)		
		
		title_page_layout = prs.slide_layouts[0] #Creates
		title_slide = prs.slides.add_slide(title_page_layout)
		title_page_title = title_slide.shapes.title
		title_page_subtitle = title_slide.placeholders[1]
		title_page_title.text = 'College Loan Project'
		title_page_subtitle.text = 'Programmed by Heyaw Meteke'
	
		
		###########################################################################
		#	EACH FUNCTION CREATES PART OF POWER POINT PRESENTATION		  #
		###########################################################################

		self.add_image_page()
		self.add_info_page()
		self.add_cost_page()
		self.export()

	
	def export (self): # Saves and exports file to folder
		name = start.collegeName
		#file_name = '~/Desktop/College_Loan_Presentations/' + start.collegeDir + '/' + 'slide.pptx'
		#prs.save(os.path.join(file_name))
		prs.save('slides.pptx')		

	def add_image_page(self):
		image_page_layout = prs.slide_layouts[6]
		image_slide =prs.slides.add_slide(image_page_layout)
		
		top = Inches(1)
		height = Inches(5.5)
		left = Inches(2)
	
		img = image_slide.shapes.add_picture(start.collegeImageDir, left, top, height)
	def add_info_page (self):
		
		name = start.collegeName
		location = start.collegeLocation
		desc = start.collegeDesc

		info_page_layout = prs.slide_layouts[1]
		info_slide = prs.slides.add_slide(info_page_layout)
		modules = info_slide.shapes
	
		title_info_page = modules.title
		body_info_page = modules.placeholders[1]
		
		title_info_page.text = name
		
		textbox = body_info_page.text_frame
		textbox.text = 'Location: ' + location
		
		p = textbox.add_paragraph()
		p.text = desc
		p.font.size = Pt(15)
	
	def add_cost_page (self):

		simple_loan = loan.simple
		tuition = start.collegeTuition

		table_page_layout = prs.slide_layouts[5]
		simple_slide = prs.slides.add_slide(table_page_layout)
		modules = compound_slide.shapes
		
		modules.title.text = 'Simple Interest'
		
		rows = 3
		cols = 3
		left = top = Inches(2.0)
		width = Inches(6.0)
		height = Inches(0.8)

		table = modules.add_table(rows, cols, left, top, width, height).table
		
		table.cell(0, 0).text = 'Year' 
		table.cell(0, 1).text = 'Interest'
		table.cell(0, 2).text = 'Balance'
		for x in range(0,11):

			if year == 1 or year == 5 or year == 10:
				table.cell(x, 0).text = str(x)
				table.cell(x, 1).text = str(simple_loan[x])
				table.cell(x, 2).text = str(simple_loan[x]+start.collegeTuition)	
				
class college_scrapper (object):
	collegeName = ''
	collegeImageDir = ''
	collegeLocation = ''
	collegeTuition = 0
	collegeDesc = ''' '''
	collegDir = ''
	loan_table = {}
	url = ''
	def __init__ (self):
		os.system('clear')
		
		print '\t\t\tCollege Loan Algebra Project'
		print '\t\t\tProgrammed by Heyaw Meteke'	
		
		driver.get('http://www.google.com')
		driver.set_window_size(1440, 900)
	
	
		self.college_search()
		self.get_image()

        def get_image (self):
		image_url = ''
                driver.get('https://www.yandex.com/images/')
                search_image = driver.find_element_by_name('text')
                search_image.clear()
		search_image.send_keys('logo ' + self.collegeName)
		search_image.send_keys(Keys.RETURN)

	
	#	try:
		time.sleep(4)
		print 'Searching for Image...'
		temp_image_url = driver.find_element_by_css_selector("*[class^='serp-item__link']").get_attribute('href')
		driver.get(temp_image_url)
		
		try:
			#Makes sure to close pop up if its open
			driver.find_element_by_css_selector("*[class^='popup__content']").click()
		except:
			#if no pop up, no worry - just to catch error
			pass
		try:
	
			image_url = driver.find_element_by_css_selector("*[class^='button2 button2_theme_action button2_size_m button2_type_link button2_pin_brick-clear button2_width_max sizes__download i-bem button2_js_inited']").get_attribute('href')

		except:
			print "First method didnt work, going to second one"
			image_url = driver.find_element_by_xpath("/html/body/div[5]/div[3]/div[2]/div[1]/div[2]/div[1]/a").get_attribute('href')	

		
		print 'Downloading Image...',
#		try:
		
		self.collegeImageDir =  'image.png'
		#	print self.collegeImageDir
		#	photo = open(self.collegeImageDir, 'wb')
		#	photo.write(requests.get(image_url).content)
		#	photo.close()

		photo = urllib.URLopener()
		photo.retrieve(image_url, 'image.png')
		print '\tOK'

#		except:

#			print 'Error downloading image'
#			print 'Placeholder is placed.  You can change it when editing file.'
#			
#			### THIS LINE DOESNT WORK
#			#self.collegeImageDir = '~/Programs/College-Loan-Web-Scrapper/placeholder.png'
#			#### DIRECTORY NOT FOUND ^ FIX !!!!
#
#			### REPLACEMENT FOR TIME BEING
#			self.collegeImageDir = 'placeholder.png'		

	def college_search (self):

		college_input = raw_input('College: ')
		college_input += ' college data'
		print 'Loading...',

		time.sleep(2)
		search_college = driver.find_element_by_name('q')
		search_college.clear()
		search_college.send_keys(college_input)
		search_college.send_keys(Keys.RETURN)
		print '\tOK'
		print 'Searching and collecting info...',
		
		time.sleep(3)
		
		try:
			driver.find_element_by_partial_link_text('CollegeData').click()

		except:
			print 'Taking longer than usual... Slow internet?(1)'
			time.sleep(3)
			driver.get('www.google.com')
			search_college = driver.find_element_by_name('q')
        	        search_college.clear()
                	search_college.send_keys(college_input)
                	search_college.send_keys(Keys.RETURN)

			driver.find_element_by_partial_link_text('CollegeData').click()
		time.sleep(5)		


		##################################################################
		#           		FIND COLLEGE INFORMATION 		 #
		##################################################################
	#	for x in range(0,5):
	#		try:
			
		self.collegeName = driver.find_element_by_xpath("//*[@id='collprofile']/div[6]/div[4]/div[2]/div[1]/h1").text	# NAME OF COLLEGE
	#		self.collegeName = driver.find_element_by_css_selector("*[class^='mainsidecontainer']").get_attribute('h1').text		

		self.collegeLocation = driver.find_element_by_css_selector("*[class^='citystate']").text
		
		self.collegeTuition = driver.find_element_by_xpath("//*[@id='section1']/table/tbody/tr[2]/td").text # TUITION COST FOR COLLEGE

		self.collegeDesc = driver.find_element_by_xpath("//*[@id='cont_overview']/p").text #BRIEF DESCRIPTION OF COLLEGE
		
		self.collegeDesc = driver.find_element_by_css_selector("*[class^='overviewtext']").text

		#	except:
		#		print 'Error!  Trying again!'
		#		time.sleep(2)
		#		if x == 4:
		#			print 'COULDNT FIND INFORMATION.  PLEASE REPORT!'
		#			exit(1)
		#		continue
		##	else:
		#		print 'Successfully found the information'
		#		break
		
		#################################################################
		print '\tOK'

		######### SEPERATE DATA AND SEND TUITION TO FIND_LOAN FUNCTION #####################

		self.collegeDir = self.collegeName.replace(' ', "_")

		self.collegeTuition = self.collegeTuition.split('$')

		tuition = ''

		tuitionTemp = self.collegeTuition[1]

		tuitionTemp = tuitionTemp.replace(' ', "")

		for char in tuitionTemp:
			try:
				works = int(char)
				tuition+= char
			except:
				pass
		self.collegeTuition = int(tuition)
		self.loan_table = find_loan(self.collegeTuition, .04)



#driver = webdriver.Firefox() # FOR VISUAL AND DEBUGGING
driver = webdriver.PhantomJS() # FOR FINAL VERSION

start = college_scrapper()

loan = loan()

prs = Presentation()
p = powerpoint()
